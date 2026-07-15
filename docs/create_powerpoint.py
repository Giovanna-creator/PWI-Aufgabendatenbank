from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).parent / "Aufgabendatenbank_Praesentation.pptx"
W, H = 13.333, 7.5
NAVY = "0B1628"
NAVY_2 = "142238"
GREEN = "35C58A"
BLUE = "4F7DF3"
ORANGE = "F5A524"
RED = "E96A6A"
WHITE = "FFFFFF"
INK = "182230"
MUTED = "607086"
LIGHT = "F4F7FA"
LINE = "DCE3EA"


def rgb(value):
    return RGBColor.from_string(value)


def rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1)
    return shape


def text(slide, value, x, y, w, h, size=18, color=INK, bold=False,
         align=PP_ALIGN.LEFT, font="Aptos", margin=0.04, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = value
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    return box


def bullet_list(slide, items, x, y, w, h, size=17, color=INK, accent=GREEN):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.03)
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(10)
        p.level = 0
        p.text = "●  " + p.text
        p.runs[0].font.color.rgb = rgb(accent)
    return box


def base_slide(prs, title, kicker=None, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = NAVY if dark else LIGHT
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(bg)
    if kicker:
        text(slide, kicker.upper(), 0.7, 0.35, 5.0, 0.3, 10,
             GREEN if dark else BLUE, True)
    text(slide, title, 0.7, 0.72, 11.8, 0.62, 28,
         WHITE if dark else NAVY, True)
    if not dark:
        rect(slide, 0.7, 1.42, 1.0, 0.04, GREEN, GREEN, False)
    return slide


def footer(slide, number, dark=False):
    color = "9EABC0" if dark else MUTED
    text(slide, "PWI · Gruppe 2 · Aufgabendatenbank", 0.7, 7.15, 5.5, 0.2, 9, color)
    text(slide, f"{number:02d}", 12.0, 7.12, 0.6, 0.22, 9, color, True, PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, title_value, body, accent=GREEN, number=None):
    rect(slide, x, y, w, h, WHITE, LINE)
    if number is not None:
        rect(slide, x + 0.25, y + 0.24, 0.5, 0.5, accent, accent)
        text(slide, str(number), x + 0.25, y + 0.28, 0.5, 0.3, 13, WHITE, True,
             PP_ALIGN.CENTER)
        tx = x + 0.88
    else:
        rect(slide, x + 0.25, y + 0.3, 0.08, 0.52, accent, accent, False)
        tx = x + 0.48
    text(slide, title_value, tx, y + 0.27, w - (tx - x) - 0.2, 0.38, 17, NAVY, True)
    text(slide, body, x + 0.25, y + 0.9, w - 0.5, h - 1.05, 13, MUTED)


def notes(slide, value):
    slide.notes_slide.notes_text_frame.text = value


prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
prs.core_properties.title = "Aufgabendatenbank – Projektpräsentation"
prs.core_properties.subject = "Praktikum Wirtschaftsinformatik, SS 2026"
prs.core_properties.author = "Joelle Giovanna Kamwa Mokam, Pharel, Danylo"

# 1 — Titel
s = base_slide(prs, "Aufgabendatenbank", "Praktikum Wirtschaftsinformatik · SS 2026", True)
text(s, "Vom abstrakten Datenschema zum nutzbaren Autorenwerkzeug",
     0.72, 1.65, 7.4, 0.75, 22, WHITE)
text(s, "Vue · Spring Boot · PostgreSQL", 0.72, 2.65, 5.0, 0.35, 13, GREEN, True)
for i, (label, x, y, col) in enumerate([
    ("ITEM", 9.0, 1.75, BLUE), ("CONTENT", 10.65, 2.75, GREEN),
    ("TAGS", 8.7, 4.1, ORANGE), ("COLLECTION", 10.55, 5.05, BLUE)
]):
    rect(s, x, y, 1.7, 0.66, NAVY_2, col)
    text(s, label, x, y + 0.19, 1.7, 0.24, 11, WHITE, True, PP_ALIGN.CENTER)
text(s, "Joelle Giovanna Kamwa Mokam · Pharel · Danylo",
     0.72, 5.85, 8.0, 0.35, 14, WHITE, True)
text(s, "Betreuung: Prof. Dr. Markus Siepermann · Johannes Kunz",
     0.72, 6.3, 8.0, 0.3, 11, "A9B5C7")
footer(s, 1, True)
notes(s, "Unser Ausgangspunkt war ein abgestimmtes Datenschema, aber noch kein "
          "nutzbares Produkt. Wir zeigen heute, wie wir daraus ein Autorenwerkzeug "
          "für Lehrende entwickelt haben.")

# 2 — Ausgangslage
s = base_slide(prs, "Ausgangslage und Projektauftrag", "01 · Problem")
card(s, 0.7, 1.75, 3.65, 3.9, "Vorgegeben", "Relationales Schema mit vielen "
     "Entitäten, Join-Tabellen und abstrakten Beziehungen.", BLUE, 1)
card(s, 4.82, 1.75, 3.65, 3.9, "Herausforderung", "Collections, Varianten und "
     "Inhaltsbausteine mussten fachlich verstanden und bedienbar werden.", ORANGE, 2)
card(s, 8.94, 1.75, 3.65, 3.9, "Unser Ergebnis", "Durchgängige Webanwendung für "
     "Erstellung, Strukturierung, Kombination und Suche von Aufgaben.", GREEN, 3)
rect(s, 2.1, 6.05, 9.1, 0.55, NAVY, NAVY)
text(s, "Nicht nur speichern – Beziehungen verständlich nutzbar machen.",
     2.1, 6.2, 9.1, 0.25, 15, WHITE, True, PP_ALIGN.CENTER)
footer(s, 2)
notes(s, "Die Aufgabe war ausdrücklich mehr als CRUD. Wir mussten entscheiden, "
          "wie technische Beziehungen als verständliche Arbeitsabläufe erscheinen.")

# 3 — Zielbild
s = base_slide(prs, "Was das System leistet", "02 · Zielbild")
features = [
    ("Aufgaben", "Metadaten und mehrere Inhalte verwalten", BLUE),
    ("Collections", "Sammlungen und geordnete Sequenzen", GREEN),
    ("Varianten", "Aufgabenfamilien horizontal abbilden", ORANGE),
    ("Metadaten", "Autoren, Lizenzen, Typen und Tags", BLUE),
    ("Darstellung", "XML-Templates mit Vorschau", GREEN),
    ("Suche", "Inhalt, Autor, Typ und Tags kombinieren", ORANGE),
]
for i, (a, b, c) in enumerate(features):
    x = 0.7 + (i % 3) * 4.12
    y = 1.75 + (i // 3) * 2.18
    card(s, x, y, 3.65, 1.75, a, b, c)
text(s, "Zielgruppe: Lehrende und Autoren von Übungsaufgaben",
     0.75, 6.35, 6.8, 0.3, 14, NAVY, True)
text(s, "Abgrenzung: kein Lern-, Abgabe- oder Korrektursystem",
     7.1, 6.35, 5.45, 0.3, 13, MUTED, False, PP_ALIGN.RIGHT)
footer(s, 3)
notes(s, "Im Mittelpunkt stehen Lehrende. Die Anwendung verwaltet Aufgaben und "
          "ihre Beziehungen; Ausführung und automatische Korrektur liegen außerhalb "
          "unseres Systems.")

# 4 — Kernmodell
s = base_slide(prs, "Das fachliche Kernmodell", "03 · Datenmodell")
rect(s, 4.77, 2.6, 3.0, 1.25, NAVY, NAVY)
text(s, "ITEM", 4.77, 2.86, 3.0, 0.34, 24, WHITE, True, PP_ALIGN.CENTER)
text(s, "Metadaten der Aufgabe", 4.77, 3.28, 3.0, 0.25, 11, "B9C5D6", False, PP_ALIGN.CENTER)
model_cards = [
    (0.8, 1.7, "ITEM CONTENT", "Text · JSON · Bild · PDF\nRolle über purpose", BLUE),
    (8.97, 1.7, "COLLECTION", "Sammlung oder Sequenz\nPosition in Join-Tabelle", GREEN),
    (0.8, 4.55, "TAGS & VALIDATOREN", "Hierarchie und Regeln\nwiederverwendbar", ORANGE),
    (8.97, 4.55, "VARIANTE", "Self-Reference über\nroot_item_id", BLUE),
]
for x, y, a, b, c in model_cards:
    card(s, x, y, 3.55, 1.6, a, b, c)
for x1, y1, x2, y2 in [(4.35, 2.5, 4.77, 2.8), (7.77, 2.8, 8.97, 2.5),
                       (4.35, 5.0, 4.77, 3.65), (7.77, 3.65, 8.97, 5.0)]:
    line = s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb("A8B4C4")
    line.line.width = Pt(2)
rect(s, 3.15, 6.35, 7.05, 0.48, "E8F7F1", "B9E8D6")
text(s, "Collection-Zugehörigkeit und Variantenbezug bleiben unabhängig.",
     3.15, 6.48, 7.05, 0.22, 13, "167A58", True, PP_ALIGN.CENTER)
footer(s, 4)
notes(s, "Das Item ist die zentrale Einheit. Inhalt wird getrennt gespeichert und "
          "über purpose verbunden. Collections organisieren verschiedene Aufgaben. "
          "root_item_id verwenden wir ausschließlich für Varianten.")

# 5 — Architektur
s = base_slide(prs, "Architektur: klar getrennte Verantwortungen", "04 · Systemdesign")
layers = [
    ("Vue 3 + Vuetify", "Komponenten, Editor, Baum, Filter", BLUE),
    ("Pinia + ApiAdapter", "Zustand, DTO-Mapping, Dummy/HTTP", GREEN),
    ("Spring Boot", "Controller → Service → Repository", ORANGE),
    ("PostgreSQL 16", "UUID · JSONB · BYTEA · Constraints", NAVY),
]
for i, (a, b, c) in enumerate(layers):
    y = 1.65 + i * 1.16
    rect(s, 1.0, y, 7.55, 0.82, WHITE if i < 3 else NAVY, c)
    text(s, a, 1.28, y + 0.19, 2.3, 0.3, 17, WHITE if i == 3 else NAVY, True)
    text(s, b, 3.65, y + 0.2, 4.55, 0.27, 14, "C6D1E0" if i == 3 else MUTED)
    if i < 3:
        text(s, "↓", 4.5, y + 0.79, 0.5, 0.4, 18, GREEN, True, PP_ALIGN.CENTER)
card(s, 9.05, 1.65, 3.25, 1.35, "Warum diese Architektur?", "Testbar, wartbar und "
     "für reale API sowie lokalen Dummy-Betrieb geeignet.", GREEN)
card(s, 9.05, 3.3, 3.25, 1.35, "Warum PostgreSQL?", "Das Schema nutzt JSONB, "
     "Binärdaten und UUIDs nativ.", BLUE)
card(s, 9.05, 4.95, 3.25, 1.35, "Warum Spring Boot?", "Transaktionen, JPA und "
     "klare Schichten – plus vorhandene Erfahrung.", ORANGE)
footer(s, 5)
notes(s, "Der Store hängt nicht direkt von Axios ab, sondern vom ApiAdapter. "
          "Im Backend liegen HTTP-Vertrag, Geschäftsregeln und Persistenz in "
          "getrennten Schichten.")

# 6 — Durchgängiger Workflow
s = base_slide(prs, "Ein Vorgang durch alle Schichten", "05 · Beispiel: Aufgabe erstellen")
steps = [
    ("1", "Dialog", "Eingaben prüfen"),
    ("2", "POST /items", "Item + UUID"),
    ("3", "POST /contents", "Inhalt speichern"),
    ("4", "item_contents", "purpose verbinden"),
    ("5", "UI", "Baum aktualisieren"),
]
for i, (n, a, b) in enumerate(steps):
    x = 0.65 + i * 2.52
    rect(s, x, 2.15, 2.0, 1.55, WHITE, LINE)
    rect(s, x + 0.64, 1.72, 0.72, 0.72, GREEN, GREEN)
    text(s, n, x + 0.64, 1.91, 0.72, 0.27, 15, WHITE, True, PP_ALIGN.CENTER)
    text(s, a, x + 0.15, 2.6, 1.7, 0.28, 16, NAVY, True, PP_ALIGN.CENTER)
    text(s, b, x + 0.15, 3.06, 1.7, 0.25, 12, MUTED, False, PP_ALIGN.CENTER)
    if i < 4:
        text(s, "→", x + 2.04, 2.65, 0.45, 0.35, 20, BLUE, True, PP_ALIGN.CENTER)
rect(s, 1.25, 4.65, 10.85, 1.25, "FFF8E9", "F1D69C")
text(s, "Teilfehler?", 1.55, 4.96, 1.55, 0.3, 17, "9A6610", True)
text(s, "Das Frontend lädt den betroffenen Bereich erneut – keine dauerhaften Ghost-Items.",
     3.15, 4.96, 8.45, 0.32, 15, INK)
text(s, "Geschäftsregeln und Referenzprüfung liegen im Backend, nicht in der Oberfläche.",
     1.55, 5.48, 9.8, 0.25, 12, MUTED)
footer(s, 6)
notes(s, "Die Item-ID muss zuerst existieren, bevor Content verknüpft werden kann. "
          "Mehrstufige Vorgänge aktualisieren die Oberfläche optimistisch. Bei einem "
          "Teilfehler synchronisieren wir gezielt mit dem Backend.")

# 7 — Collections
s = base_slide(prs, "Strukturieren: Collections und Varianten", "06 · Fachliche Logik")
card(s, 0.7, 1.72, 3.8, 3.8, "Ungeordnete Collection",
     "Thematische Gruppierung\n\nposition = NULL\n\nReihenfolge ohne Bedeutung", GREEN)
card(s, 4.77, 1.72, 3.8, 3.8, "Geordnete Collection",
     "Lernsequenz\n\nposition = 1, 2, 3 …\n\nBackend normalisiert Positionen", BLUE)
card(s, 8.84, 1.72, 3.8, 3.8, "Horizontale Variante",
     "Gleiche Aufgabenfamilie\n\nroot_item_id\n\nUnabhängig von Collections", ORANGE)
rect(s, 1.55, 6.0, 10.2, 0.56, NAVY, NAVY)
text(s, "Drag-and-Drop  →  Service-Operation  →  konsistente Geschwisterpositionen",
     1.55, 6.16, 10.2, 0.25, 14, WHITE, True, PP_ALIGN.CENTER)
footer(s, 7)
notes(s, "Die Reihenfolge gehört zur Mitgliedschaft, nicht zum Item. Deshalb "
          "berechnet das Backend die endgültigen Positionen. Ein Verschieben darf "
          "rootItemId niemals verändern.")

# 8 — Inhalte & Metadaten
s = base_slide(prs, "Inhalte, Metadaten und Wiederfinden", "07 · Autorenworkflow")
columns = [
    ("INHALTE", ["mehrere Blöcke", "Text, Bild, PDF", "eigener purpose"], BLUE),
    ("DARSTELLUNG", ["XML-Template", "Purpose-Reihenfolge", "Live-Vorschau"], GREEN),
    ("TAGS & SUCHE", ["hierarchischer Baum", "Autor und Typ", "kombinierte Filter"], ORANGE),
]
for i, (a, items, c) in enumerate(columns):
    x = 0.7 + i * 4.12
    rect(s, x, 1.75, 3.65, 4.55, WHITE, LINE)
    rect(s, x, 1.75, 3.65, 0.75, c, c)
    text(s, a, x + 0.25, 1.98, 3.15, 0.25, 15, WHITE, True, PP_ALIGN.CENTER)
    bullet_list(s, items, x + 0.35, 2.9, 2.95, 2.2, 16, INK, c)
    text(s, ["JSONB / BYTEA", "CodeMirror 6", "JPA Specifications"][i],
         x + 0.35, 5.55, 2.95, 0.28, 12, c, True, PP_ALIGN.CENTER)
text(s, "Speichern allein genügt nicht: Aufgaben müssen verständlich dargestellt und wiedergefunden werden.",
     1.2, 6.55, 10.9, 0.3, 14, NAVY, True, PP_ALIGN.CENTER)
footer(s, 8)
notes(s, "Die Aufgabe ist kein großes Textfeld. Sie besteht aus mehreren Blöcken. "
          "Templates ordnen diese Blöcke, und hierarchische Tags plus kombinierte "
          "Filter machen sie wieder auffindbar.")

# 9 — Qualität
s = base_slide(prs, "Qualität und Robustheit", "08 · Absicherung")
for value, label, x, c in [
    ("107", "Frontend-Tests", 0.8, BLUE),
    ("92", "Backend-Tests", 3.7, GREEN),
    ("✓", "Typecheck & Build", 6.6, ORANGE),
    ("SHA", "reproduzierbare Images", 9.5, NAVY),
]:
    rect(s, x, 1.72, 2.55, 1.72, WHITE, c)
    text(s, value, x, 1.98, 2.55, 0.55, 30, c, True, PP_ALIGN.CENTER)
    text(s, label, x + 0.15, 2.75, 2.25, 0.25, 12, MUTED, True, PP_ALIGN.CENTER)
card(s, 0.8, 3.95, 3.55, 1.65, "Geschäftslogik", "Collections, Positionen, "
     "Varianten und Referenzen", BLUE)
card(s, 4.88, 3.95, 3.55, 1.65, "Fehlerpfade", "Resynchronisation nach "
     "partiellen API-Fehlern", GREEN)
card(s, 8.96, 3.95, 3.55, 1.65, "Integration", "Feature-Branches, Pull Requests "
     "und Reviews", ORANGE)
text(s, "CI erstellt nach Merge in main Backend- und Frontend-Images in GHCR.",
     1.3, 6.25, 10.7, 0.3, 14, NAVY, True, PP_ALIGN.CENTER)
footer(s, 9)
notes(s, "Wir haben nicht nur Happy Paths getestet. Kritisch waren Positionslogik, "
          "Teilfehlerszenarien und die Trennung zwischen Collection und Variante.")

# 10 — Organisation
s = base_slide(prs, "Projektorganisation", "09 · Zusammenarbeit")
timeline = [("APR/MAI", "Analyse"), ("MAI", "Grundaufbau"), ("JUNI", "Kernfunktionen"),
            ("JUN/JUL", "Erweiterungen"), ("JULI", "Stabilisierung")]
for i, (date, phase) in enumerate(timeline):
    x = 0.8 + i * 2.47
    rect(s, x, 1.82, 2.0, 0.08, GREEN, GREEN, False)
    rect(s, x, 1.62, 0.48, 0.48, NAVY, NAVY)
    text(s, str(i + 1), x, 1.74, 0.48, 0.2, 11, WHITE, True, PP_ALIGN.CENTER)
    text(s, date, x, 2.18, 2.0, 0.22, 10, BLUE, True)
    text(s, phase, x, 2.52, 2.0, 0.25, 14, NAVY, True)
card(s, 0.8, 3.45, 3.55, 2.05, "GitHub Project", "Arbeitspakete über "
     "In progress, In review und Closed/Done nachverfolgt.", BLUE)
card(s, 4.88, 3.45, 3.55, 2.05, "Kommunikation", "Ein bis zwei Meetings pro "
     "Woche – abhängig von Fortschritt und Klärungsbedarf.", GREEN)
card(s, 8.96, 3.45, 3.55, 2.05, "Zusammenarbeit", "Pull-Request-Reviews und "
     "Pair Programming bei komplexen Integrationsfragen.", ORANGE)
text(s, "Joelle: Leitung + Full Stack/Deployment  ·  Danylo: Frontend/Backend  ·  Pharel: Backend/Tests",
     0.9, 6.2, 11.55, 0.3, 12, MUTED, True, PP_ALIGN.CENTER)
footer(s, 10)
notes(s, "GitHub war unser Arbeitsboard und nicht nur ein Codearchiv. Durch die "
          "Reduktion von fünf auf drei aktive Personen mussten wir neu priorisieren "
          "und schichtübergreifend arbeiten.")

# 11 — Grenzen
s = base_slide(prs, "Grenzen und bewusste Entscheidungen", "10 · Transparenz")
card(s, 0.75, 1.75, 3.7, 4.25, "Im Kern umgesetzt",
     "Items und Contents\nCollections und Varianten\nTags und Suche\nValidator-Zuordnung\nTemplates", GREEN)
card(s, 4.82, 1.75, 3.7, 4.25, "Teilweise vorbereitet",
     "Modifier-Verwaltung\nContent-Tags in der UI\nTyp-Kompatibilitätsfilter\nContent-Wiederverwendung", ORANGE)
card(s, 8.89, 1.75, 3.7, 4.25, "Externe Abhängigkeiten",
     "zentrale Authentifizierung\nValidator-Ausführung\nArgoCD-Zielplattform\nproduktives Cluster", BLUE)
text(s, "Container und Kubernetes-Konfiguration sind vorhanden; die Zielinfrastruktur wurde nicht bereitgestellt.",
     1.05, 6.35, 11.2, 0.32, 13, NAVY, True, PP_ALIGN.CENTER)
footer(s, 11)
notes(s, "Wir trennen klar zwischen umgesetzt, vorbereitet und extern abhängig. "
          "Das produktive Deployment scheiterte nicht am Container-Build, sondern "
          "an der fehlenden Zielplattform.")

# 12 — Ergebnis
s = base_slide(prs, "Ergebnis", "11 · Fazit", True)
text(s, "Ein abstraktes Schema wurde zu einem bedienbaren Autorenwerkzeug.",
     0.75, 1.62, 11.7, 0.58, 24, WHITE, True, PP_ALIGN.CENTER)
for i, (a, b, c) in enumerate([
    ("DURCHGÄNGIG", "Datenbank, Backend und Frontend integriert", BLUE),
    ("KONSISTENT", "Collections, Positionen und Varianten sauber getrennt", GREEN),
    ("BELASTBAR", "getestet, containerisiert und erweiterbar", ORANGE),
]):
    x = 0.85 + i * 4.15
    rect(s, x, 2.7, 3.55, 2.45, NAVY_2, c)
    text(s, a, x + 0.25, 3.05, 3.05, 0.3, 15, c, True, PP_ALIGN.CENTER)
    text(s, b, x + 0.35, 3.7, 2.85, 0.7, 16, WHITE, False, PP_ALIGN.CENTER)
text(s, "Die Stärke liegt nicht in einzelnen CRUD-Masken, sondern in den verständlichen Beziehungen.",
     1.15, 5.85, 11.0, 0.45, 15, "C7D1E0", False, PP_ALIGN.CENTER)
footer(s, 12, True)
notes(s, "Unser wichtigstes Ergebnis ist die Übersetzung der Beziehungen in "
          "nachvollziehbare Arbeitsabläufe. Darauf kann eine nächste Iteration "
          "mit Authentifizierung und externen Diensten aufbauen.")

# 13 — Demo
s = base_slide(prs, "Live-Demo", "12 · Anwendung", True)
demo = [
    "Aufgabe mit Autor, Lizenz und Typ erstellen",
    "Content hinzufügen und Template-Vorschau zeigen",
    "Tag zuweisen und über Filter wiederfinden",
    "Geordnete Collection erstellen und verschieben",
    "Variante und Validator zeigen",
    "Neu laden: Persistenz nachweisen",
]
for i, item in enumerate(demo):
    y = 1.55 + i * 0.78
    rect(s, 1.45, y, 0.48, 0.48, GREEN, GREEN)
    text(s, str(i + 1), 1.45, y + 0.12, 0.48, 0.2, 11, WHITE, True, PP_ALIGN.CENTER)
    text(s, item, 2.15, y + 0.08, 8.9, 0.3, 17, WHITE)
text(s, "Vorbereitete Daten · keine Infrastrukturänderungen während der Demo",
     2.15, 6.45, 8.9, 0.28, 12, "A8B5C8")
footer(s, 13, True)
notes(s, "Die Demo beginnt mit laufendem System und vorbereiteten Daten. Wir zeigen "
          "einen zusammenhängenden Autorenworkflow und vermeiden redundante CRUD-Dialoge.")

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(OUT)
