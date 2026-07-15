from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


TEMPLATE = Path(r"C:\Users\hp\Downloads\Zwischenpraesentation-Aufgabendatenbank1.pptx")
OUT = Path(__file__).parents[1] / "Projektpraesentation_final.pptx"
NAVY, TEAL, ORANGE = "1E2761", "00A98F", "FFB547"
WHITE, LIGHT, INK, GRAY, LINE = "FFFFFF", "F5F7FA", "202533", "667085", "D9DEE8"
W, H = 13.333, 7.5


def color(value):
    return RGBColor.from_string(value)


def box(slide, x, y, w, h, fill=WHITE, stroke=LINE, rounded=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.color.rgb = color(stroke)
    shape.line.width = Pt(1)
    return shape


def txt(slide, value, x, y, w, h, size=20, fill=INK, bold=False,
        align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font="Calibri"):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = tf.margin_right = Inches(.04)
    tf.margin_top = tf.margin_bottom = Inches(.03)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text, p.alignment = value, align
    p.font.name, p.font.size = font, Pt(size)
    p.font.bold, p.font.color.rgb = bold, color(fill)
    return shape


def bullets(slide, items, x, y, w, h, size=18, fill=INK):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + item
        p.font.name, p.font.size = "Calibri", Pt(size)
        p.font.color.rgb = color(fill)
        p.space_after = Pt(10)
    return shape


def clear_slides(prs):
    for slide_id in list(prs.slides._sldIdLst):
        prs.part.drop_rel(slide_id.rId)
        prs.slides._sldIdLst.remove(slide_id)


def slide_base(prs, title, section, speaker, number, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color(NAVY if dark else LIGHT)
    txt(slide, section.upper(), .65, .22, 5.0, .22, 11, TEAL, True)
    txt(slide, title, .65, .58, 9.6, .52, 30, WHITE if dark else NAVY, True)
    if not dark:
        box(slide, .65, 1.23, 1.05, .045, TEAL, TEAL, False)
    presenter(slide, speaker, dark)
    txt(slide, "PWI · Gruppe 2", .65, 7.14, 2.4, .18, 10, "B7C0D0" if dark else GRAY)
    txt(slide, f"{number:02d}", 12.0, 7.12, .65, .2, 10,
        "B7C0D0" if dark else GRAY, True, PP_ALIGN.RIGHT)
    return slide


def presenter(slide, name, dark=False):
    pill = box(slide, 10.45, .28, 2.25, .55, NAVY if not dark else "2A3570",
               NAVY if not dark else "46518A")
    initials = "".join(part[0] for part in name.split()[:2]).upper()
    box(slide, 10.58, .38, .34, .34, TEAL, TEAL)
    txt(slide, initials, 10.58, .45, .34, .14, 9, WHITE, True, PP_ALIGN.CENTER)
    txt(slide, "PRÄSENTIERT VON", 11.02, .36, 1.5, .13, 8, "BEC8DA")
    txt(slide, name, 11.02, .52, 1.5, .18, 11, WHITE, True)


def card(slide, x, y, w, h, heading, body, accent=TEAL):
    box(slide, x, y, w, h, WHITE, LINE)
    box(slide, x + .22, y + .24, .07, .5, accent, accent, False)
    txt(slide, heading, x + .43, y + .23, w - .62, .32, 19, NAVY, True)
    txt(slide, body, x + .24, y + .74, w - .48, h - .84, 16, GRAY)


def placeholder(slide, x, y, w, h, label, hint):
    shape = box(slide, x, y, w, h, "EEF1F5", "AEB7C6")
    shape.line.width = Pt(1.5)
    txt(slide, "▧", x, y + h * .25, w, .5, 30, TEAL, True, PP_ALIGN.CENTER)
    txt(slide, label, x + .3, y + h * .51, w - .6, .3, 17, NAVY, True, PP_ALIGN.CENTER)
    txt(slide, hint, x + .35, y + h * .68, w - .7, .38, 12, GRAY, False, PP_ALIGN.CENTER)


def demo_bar(slide, title, actions, seconds="60 s"):
    box(slide, .72, 6.22, 11.9, .62, NAVY, NAVY)
    box(slide, .9, 6.34, 1.48, .35, TEAL, TEAL)
    txt(slide, f"MINI-DEMO · {seconds}", .9, 6.43, 1.48, .14, 9, WHITE, True, PP_ALIGN.CENTER)
    txt(slide, title, 2.58, 6.34, 2.55, .22, 14, WHITE, True)
    txt(slide, "  →  ".join(actions), 5.0, 6.35, 7.25, .24, 13, WHITE)


def notes(slide, value):
    frame = slide.notes_slide.notes_text_frame
    if frame is not None:
        frame.text = value


prs = Presentation(TEMPLATE)
clear_slides(prs)
prs.core_properties.title = "Aufgabendatenbank – Abschlusspräsentation"
prs.core_properties.author = "Joelle Giovanna Kamwa Mokam, Danylo, Pharel"

# 1
s = slide_base(prs, "Aufgabendatenbank", "Abschlusspräsentation · SS 2026", "Team", 1, True)
txt(s, "Vom vorgegebenen Datenschema zum nutzbaren Autorenwerkzeug",
    .75, 1.6, 8.1, .8, 25, WHITE)
for i, label in enumerate(["DATENMODELL", "WEBANWENDUNG", "AUTORENWORKFLOW"]):
    box(s, .78 + i * 2.75, 2.8, 2.38, .62, "2A3570", "46518A")
    txt(s, label, .78 + i * 2.75, 3.0, 2.38, .2, 12, TEAL, True, PP_ALIGN.CENTER)
txt(s, "Joelle Giovanna Kamwa Mokam  ·  Danylo  ·  Pharel",
    .78, 5.55, 8.4, .32, 16, WHITE, True)
txt(s, "Betreuung: Prof. Dr. Markus Siepermann · Johannes Kunz",
    .78, 6.02, 8.4, .28, 13, "C4CCDA")
notes(s, "Wir zeigen heute nicht nur einzelne Masken, sondern wie wir ein abstraktes "
         "Schema in zusammenhängende Autorenabläufe übersetzt haben.")

# 2
s = slide_base(prs, "Projektauftrag und Ziel", "Ausgangslage", "Joelle", 2)
card(s, .72, 1.62, 3.65, 3.95, "Ausgangspunkt",
     "Ein abgestimmtes relationales Schema – aber noch keine praktisch nutzbare Anwendung.")
card(s, 4.84, 1.62, 3.65, 3.95, "Unsere Aufgabe",
     "Entitäten und Beziehungen fachlich verstehen und in klare Interaktionen übersetzen.")
card(s, 8.96, 1.62, 3.65, 3.95, "Ziel",
     "Ein Autorenwerkzeug für Lehrende: Aufgaben erstellen, strukturieren, kombinieren und wiederfinden.")
txt(s, "Kein Lern- oder Korrektursystem", 1.2, 5.82, 3.5, .25, 15, GRAY, True)
txt(s, "Fokus: Verwaltung und Modellierung", 8.25, 5.82, 4.0, .25, 15, NAVY, True, PP_ALIGN.RIGHT)
notes(s, "Die fachliche Herausforderung lag in den Beziehungen. Nutzende sollen "
         "nicht über Fremdschlüssel nachdenken, sondern über Aufgaben, Inhalte und Strukturen.")

# 3
s = slide_base(prs, "Vorgegebenes Datenbankschema", "Modellgrundlage", "Joelle", 3)
placeholder(s, .72, 1.52, 8.05, 4.95, "DATENBANKSCHEMA EINFÜGEN",
            "breite, gut lesbare PNG-Version verwenden")
card(s, 9.1, 1.52, 3.5, 1.35, "Kern", "Item und ItemContent")
card(s, 9.1, 3.08, 3.5, 1.35, "Organisation", "Collections und Positionen")
card(s, 9.1, 4.64, 3.5, 1.35, "Metadaten", "Tags, Typen, Lizenzen, Regeln")
notes(s, "Hier wird das Originalschema gezeigt. Wir markieren beim Sprechen nur drei "
         "Bereiche: Kernobjekte, Organisation und Metadaten. Nicht jede Tabelle einzeln vorlesen.")

# 4
s = slide_base(prs, "Architektur und Technologieentscheidungen", "Systemdesign", "Joelle", 4)
layers = [
    ("FRONTEND", "Vue 3 · TypeScript · Vuetify · Pinia", "Interaktion und Zustand"),
    ("REST", "Axios · typisierte DTOs · ApiAdapter", "stabiler Vertrag"),
    ("BACKEND", "Java 21 · Spring Boot · JPA", "Regeln und Transaktionen"),
    ("DATENBANK", "PostgreSQL 16 · UUID · JSONB · BYTEA", "Integrität und Persistenz"),
]
for i, (a, b, c) in enumerate(layers):
    y = 1.5 + i * 1.14
    box(s, .82, y, 8.15, .78, WHITE if i < 3 else NAVY, TEAL if i == 0 else LINE)
    txt(s, a, 1.08, y + .21, 1.42, .22, 14, TEAL, True)
    txt(s, b, 2.7, y + .18, 3.85, .26, 17, WHITE if i == 3 else NAVY, True)
    txt(s, c, 6.6, y + .2, 2.05, .22, 13, "C9D0DE" if i == 3 else GRAY, False, PP_ALIGN.RIGHT)
card(s, 9.36, 1.5, 3.1, 1.45, "Warum passend?", "Vorwissen im Team und Anschluss an das bereitgestellte Vue-Frontend.")
card(s, 9.36, 3.2, 3.1, 1.45, "Warum PostgreSQL?", "Das Schema benötigt JSONB, Binärdaten und UUIDs nativ.")
card(s, 9.36, 4.9, 3.1, 1.15, "Betrieb", "Docker Compose und GHCR-Images.")
notes(s, "Die Technologien wurden nicht nur nach Bekanntheit gewählt. Sie passen zu "
         "Schema, Plattformgerüst und benötigten Transaktionen.")

# 5
s = slide_base(prs, "Item, Content und Referenzdaten", "Entitätsgruppe 1", "Danylo", 5)
card(s, .72, 1.5, 3.48, 1.35, "Item", "Metadaten: Autor, Lizenz, Typ und optionales Template.")
card(s, .72, 3.02, 3.48, 1.35, "ItemContent", "Text/JSON oder Datei; Verbindung zum Item über purpose.")
card(s, .72, 4.54, 3.48, 1.35, "Referenzen", "Author · License · ItemType · ContentType")
placeholder(s, 4.55, 1.5, 8.05, 4.35, "SCREENSHOT: AUFGABE + CONTENT-EDITOR",
            "Metadaten und zwei unterschiedliche Inhaltsbausteine")
demo_bar(s, "Aufgabe erstellen", ["Metadaten wählen", "Textinhalt ergänzen", "Bild/PDF zeigen"], "60 s")
notes(s, "Zuerst das Modell erklären, danach direkt in der Anwendung zeigen. "
         "Wichtig: purpose liegt auf der Verbindung, nicht im Content selbst.")

# 6
s = slide_base(prs, "Collections und Varianten", "Entitätsgruppe 2", "Danylo", 6)
card(s, .72, 1.5, 3.72, 1.48, "Ungeordnete Collection",
     "Thematische Gruppe; position bleibt NULL.")
card(s, .72, 3.16, 3.72, 1.48, "Geordnete Collection",
     "Sequenz; Backend verwaltet 1, 2, 3 …")
card(s, .72, 4.82, 3.72, 1.18, "Variante",
     "root_item_id verbindet eine Variante mit der Ausgangsaufgabe.")
placeholder(s, 4.82, 1.5, 7.78, 4.5, "SCREENSHOT: TREEVIEW + VARIANTENPANEL",
            "Collection-Knoten, Positionen und Variantenbereich")
demo_bar(s, "Struktur materialisieren", ["Collection anlegen", "ordnen per Drag-and-Drop", "Variante zeigen"], "60 s")
notes(s, "Collections bilden die vertikale Organisation. Varianten sind die horizontale "
         "Beziehung. Die aktuelle UI verwendet ein Item entweder als Collection-Träger "
         "oder als Ausgangsaufgabe für Varianten; das ist eine Anwendungsentscheidung.")

# 7
s = slide_base(prs, "Hierarchische Tags und Suche", "Entitätsgruppe 3", "Danylo", 7)
placeholder(s, .72, 1.5, 7.5, 4.5, "SCREENSHOT: TAG-BAUM UND FILTER",
            "Tag-Pfad, Chips und kombinierte Suchleiste")
card(s, 8.58, 1.5, 4.02, 1.25, "Datenbank", "Self-Reference über parent_tag_id.")
card(s, 8.58, 2.95, 4.02, 1.25, "Backend", "Flache Tag-Liste und Zuordnungsendpunkte.")
card(s, 8.58, 4.4, 4.02, 1.6, "Frontend", "Rekursiver Baum; Suche nach Text, Autor, Typ und Tag.")
demo_bar(s, "Wiederfinden", ["Tag-Pfad anlegen", "zuweisen", "Filter kombinieren"], "60 s")
notes(s, "Tags erhalten eine eigene Folie, weil Baumaufbau, Zuordnung und Filterung "
         "zusammen einen großen Workflow bilden.")

# 8
s = slide_base(prs, "Templates, Validatoren und Modifier", "Entitätsgruppe 4", "Pharel", 8)
card(s, .72, 1.5, 3.65, 2.0, "Representation Template",
     "XML legt die Reihenfolge der purpose-Blöcke in der Vorschau fest.")
card(s, .72, 3.72, 3.65, 2.0, "Validator",
     "Wiederverwendbare Regel; Speicherung und Item-Zuordnung umgesetzt.")
placeholder(s, 4.72, 1.5, 7.88, 4.22, "SCREENSHOT: TEMPLATE + VALIDATOR",
            "XML-Editor, Vorschau und zugeordnete Regel")
box(s, 9.05, 5.05, 3.2, .48, "FFF5E4", "EBCB8B")
txt(s, "Modifier: strukturell vorbereitet", 9.05, 5.2, 3.2, .18, 12, "875D12", True, PP_ALIGN.CENTER)
demo_bar(s, "Darstellung und Regeln", ["Template umsortieren", "Vorschau öffnen", "Validator zuordnen"], "60 s")
notes(s, "Validatoren werden in dieser Anwendung nicht ausgeführt. Das ist eine "
         "Systemgrenze. Modifier sind nur strukturell vorbereitet und werden offen so benannt.")

# 9
s = slide_base(prs, "Konsistenz über alle Schichten", "Backend und Datenbank", "Pharel", 9)
steps = [
    ("1", "Vue-Dialog", "Eingaben"),
    ("2", "REST-DTO", "UUIDs"),
    ("3", "Service", "Regeln"),
    ("4", "Repository", "JPA"),
    ("5", "PostgreSQL", "Constraints"),
]
for i, (n, a, b) in enumerate(steps):
    x = .6 + i * 2.53
    box(s, x, 1.72, 2.05, 1.7, WHITE, LINE)
    box(s, x + .7, 1.45, .65, .65, TEAL, TEAL)
    txt(s, n, x + .7, 1.65, .65, .18, 12, WHITE, True, PP_ALIGN.CENTER)
    txt(s, a, x + .15, 2.35, 1.75, .25, 17, NAVY, True, PP_ALIGN.CENTER)
    txt(s, b, x + .15, 2.82, 1.75, .2, 13, GRAY, False, PP_ALIGN.CENTER)
box(s, 1.15, 4.05, 11.0, 1.45, WHITE, TEAL)
txt(s, "Beispiel: geordnete Collection", 1.45, 4.35, 3.35, .3, 19, NAVY, True)
txt(s, "Backend berechnet Positionen und normalisiert Geschwister transaktional.",
    4.65, 4.34, 6.95, .32, 17, INK)
txt(s, "Bei einem Teilfehler lädt das Frontend den betroffenen Bereich erneut.",
    4.65, 4.88, 6.95, .25, 14, GRAY)
demo_bar(s, "Persistenz prüfen", ["Änderung speichern", "Seite neu laden", "gleichen Stand zeigen"], "30 s")
notes(s, "Diese Folie verbindet Frontend, Backend und Datenbank. Die wichtigste "
         "Botschaft: Der Server bleibt Autorität für fachliche Konsistenz.")

# 10
s = slide_base(prs, "Qualitätssicherung und Delivery", "Engineering", "Pharel", 10)
metrics = [("107", "Frontend-Tests"), ("92", "Backend-Tests"), ("✓", "Typecheck & Build")]
for i, (value, label) in enumerate(metrics):
    x = .75 + i * 2.72
    box(s, x, 1.55, 2.35, 1.45, WHITE, TEAL)
    txt(s, value, x, 1.82, 2.35, .48, 30, TEAL, True, PP_ALIGN.CENTER)
    txt(s, label, x, 2.48, 2.35, .22, 13, GRAY, True, PP_ALIGN.CENTER)
placeholder(s, 8.98, 1.55, 3.62, 2.95, "OPTIONAL: TESTLAUF",
            "grünes Terminal oder GitHub Check")
card(s, .75, 3.35, 3.55, 1.7, "Tests", "Store, Services, Controller und Fehlerpfade.")
card(s, 4.68, 3.35, 3.55, 1.7, "Pull Requests", "Review vor Merge; Fixes in kleinen Branches.")
card(s, .75, 5.32, 7.48, .72, "Delivery", "GitHub Actions baut Backend- und Frontend-Images mit Commit-SHA.")
demo_bar(s, "Robustheit", ["geordnet umschalten", "neu laden", "Positionen vergleichen"], "30 s")
notes(s, "Die Zahlen groß zeigen, aber kurz erklären, was getestet wurde. "
         "Nicht nur Anzahl, sondern kritische Regeln und Fehlerpfade hervorheben.")

# 11
s = slide_base(prs, "Projektorganisation und Zusammenarbeit", "Vorgehensmodell", "Joelle", 11)
placeholder(s, .72, 1.48, 5.7, 3.9, "SCREENSHOT: GITHUB PROJECT",
            "In progress · In review · Closed/Done")
placeholder(s, 6.78, 1.48, 5.82, 3.9, "SCREENSHOT: DISCORD-MEETING",
            "Teamgespräch oder Bildschirmfreigabe")
for i, (value, label) in enumerate([
    ("1–2×", "Meetings pro Woche"), ("PR", "Review vor Merge"), ("PAIR", "gemeinsame Sessions")
]):
    x = .9 + i * 4.05
    box(s, x, 5.65, 3.65, .72, WHITE, LINE)
    txt(s, value, x + .18, 5.84, .78, .23, 17, TEAL, True)
    txt(s, label, x + .98, 5.86, 2.42, .2, 14, NAVY, True)
notes(s, "GitHub war Arbeitsorganisation, nicht nur Versionsverwaltung. Meetings "
         "fanden abhängig von Fortschritt und Fragen ein- bis zweimal pro Woche statt.")

# 12
s = slide_base(prs, "Grenzen und nächste Schritte", "Transparenz", "Joelle", 12)
card(s, .72, 1.5, 3.7, 4.55, "Bewusst offen",
     "• vollständige Modifier-UI\n\n• Content-Tag-Workflow\n\n• Typ-Kompatibilitätsfilter", ORANGE)
card(s, 4.82, 1.5, 3.7, 4.55, "Externe Abhängigkeiten",
     "• zentrale Authentifizierung\n\n• Validator-Ausführung\n\n• ArgoCD-Zielplattform", ORANGE)
card(s, 8.92, 1.5, 3.7, 4.55, "Nächste Schritte",
     "• PostgreSQL-Integrationstests\n\n• Rollen und Rechte\n\n• produktives Deployment", TEAL)
txt(s, "Container und Values-Konfiguration sind vorhanden; die Zielinfrastruktur wurde nicht bereitgestellt.",
    1.05, 6.35, 11.15, .3, 14, NAVY, True, PP_ALIGN.CENTER)
notes(s, "Offene Punkte nicht verstecken, sondern sauber nach internem Scope und "
         "externer Abhängigkeit unterscheiden.")

# 13
s = slide_base(prs, "Fazit und Fragen", "Ergebnis", "Team", 13, True)
txt(s, "Aus einem abstrakten Schema wurde ein getestetes, bedienbares Autorenwerkzeug.",
    1.1, 1.65, 11.1, .75, 25, WHITE, True, PP_ALIGN.CENTER)
for i, (a, b) in enumerate([
    ("MODELLIERT", "Schemaelemente fachlich übersetzt"),
    ("INTEGRIERT", "Frontend · Backend · PostgreSQL"),
    ("ERWEITERBAR", "klare Grenzen und nächste Schritte"),
]):
    x = .85 + i * 4.15
    box(s, x, 2.95, 3.55, 2.0, "2A3570", "46518A")
    txt(s, a, x, 3.28, 3.55, .25, 15, TEAL, True, PP_ALIGN.CENTER)
    txt(s, b, x + .3, 3.9, 2.95, .55, 17, WHITE, False, PP_ALIGN.CENTER)
txt(s, "Vielen Dank.", 1.0, 5.72, 11.3, .45, 23, WHITE, True, PP_ALIGN.CENTER)
notes(s, "Die Mini-Demos wurden bereits in den fachlichen Blöcken gezeigt. Deshalb "
         "endet die Präsentation mit dem Ergebnis und geht direkt in die Fragerunde.")

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(OUT)
